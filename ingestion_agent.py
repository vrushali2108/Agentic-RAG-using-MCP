
# agents/ingestion_agent.py

import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import pandas as pd
import uuid

class IngestionAgent:
    def __init__(self):
        self.sender = "IngestionAgent"
        self.receiver = "RetrievalAgent"

    def run(self, mcp_message):
        try:
            trace_id = mcp_message.get("trace_id", str(uuid.uuid4()))
            payload = mcp_message.get("payload", {})
            file_path = payload.get("file_path", "")
            filename = os.path.basename(file_path)
            ext = filename.split('.')[-1].lower()

            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")

            if ext == 'pdf':
                text = self._parse_pdf(file_path)
            elif ext == 'docx':
                text = self._parse_docx(file_path)
            elif ext == 'pptx':
                text = self._parse_pptx(file_path)
            elif ext == 'csv':
                text = self._parse_csv(file_path)
            elif ext in ['txt', 'md']:
                text = self._parse_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")

            return {
                "sender": self.sender,
                "receiver": self.receiver,
                "type": "CONTEXT_DOCUMENT_PARSED",
                "trace_id": trace_id,
                "payload": {
                    "filename": filename,
                    "text": text
                }
            }

        except Exception as e:
            return {
                "sender": self.sender,
                "receiver": "UI",
                "type": "CONTEXT_ERROR",
                "trace_id": mcp_message.get("trace_id", str(uuid.uuid4())),
                "payload": {
                    "error": str(e)
                }
            }

    def _parse_pdf(self, path):
        doc = fitz.open(path)
        return "\n".join([page.get_text() for page in doc])

    def _parse_docx(self, path):
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _parse_pptx(self, path):
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    def _parse_csv(self, path):
        df = pd.read_csv(path)
        return df.to_string()

    def _parse_txt(self, path):
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
